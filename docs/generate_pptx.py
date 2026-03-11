"""
AgriScan Pro — パワーポイント資料生成スクリプト
概要資料 + 操作マニュアルの2つのPPTXを生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = DOCS_DIR

# ブランドカラー
BG_DARK = RGBColor(0x0A, 0x0F, 0x0D)
BG_CARD = RGBColor(0x14, 0x1E, 0x19)
BG_CARD2 = RGBColor(0x1A, 0x25, 0x20)
GREEN_PRIMARY = RGBColor(0x43, 0xA0, 0x47)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
GREEN_SOFT = RGBColor(0xA5, 0xD6, 0xA7)
TEXT_WHITE = RGBColor(0xE8, 0xF5, 0xE9)
TEXT_MUTED = RGBColor(0x8A, 0xA8, 0x94)
ACCENT_AMBER = RGBColor(0xFF, 0x98, 0x00)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, bold=False, color=TEXT_WHITE,
                 alignment=PP_ALIGN.LEFT, font_name='Meiryo'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=13, color=TEXT_WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Meiryo'
        p.space_after = Pt(6)
    return txBox


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        kwargs = {'left': Inches(left), 'top': Inches(top)}
        if width:
            kwargs['width'] = Inches(width)
        if height:
            kwargs['height'] = Inches(height)
        slide.shapes.add_picture(img_path, **kwargs)
        return True
    return False


def add_footer(slide, text="AgriScan Pro  |  ドローン画像×AI解析  |  2026"):
    add_text_box(slide, 0.5, 7.1, 9, 0.3, text,
                 font_size=8, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# =============================================
# 概要資料
# =============================================
def create_overview():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- スライド1: 表紙 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_shape_bg(slide, 0, 0, 10, 7.5, BG_DARK)

    # グリーンのアクセントライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(3.2), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_text_box(slide, 1, 1.5, 8, 1, '🌾', font_size=60, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.2, 8, 0.8, 'AgriScan Pro',
                 font_size=44, bold=True, color=GREEN_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.5, 8, 0.6, 'ドローン画像 × AI解析で圃場管理を次世代へ',
                 font_size=22, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 8, 0.4, '病害虫の早期発見 ・ 雑草密度マッピング ・ インフラ点検 ・ 散布品質評価',
                 font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.0, 8, 0.4, '対象エリア: 福岡県筑前町・朝倉エリア  |  2026年3月',
                 font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # --- スライド2: AgriScan Proとは ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, 'AgriScan Pro とは',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_shape_bg(slide, 0.5, 1.2, 9.0, 2.2, BG_CARD)
    add_text_box(slide, 0.8, 1.4, 8.5, 2.0,
                 'AgriScan Proは、ドローンで撮影した圃場画像をAI（Claude Vision API）で\n'
                 '自動解析し、病害虫の早期発見・雑草密度の把握・インフラ点検・\n'
                 '散布品質の評価を一元管理するWebアプリケーションです。\n\n'
                 '九州地方の稲作（ひのひかり・元気つくし・夢つくし）に\n'
                 '特化した解析プロンプトを搭載しています。',
                 font_size=15, color=TEXT_WHITE)

    # 4つの特徴カード
    features = [
        ('🔬', '病害虫診断', 'AI画像認識で\n11種の病害虫を検出'),
        ('🌿', '雑草マッピング', '密度ヒートマップ\n重点散布エリア提案'),
        ('🏗️', 'インフラ点検', '畦畔崩れ・水路詰まり\n獣害痕の検出'),
        ('✈️', '散布品質評価', '散布ムラの可視化\nルート改善提案'),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = 0.5 + i * 2.3
        add_shape_bg(slide, x, 3.8, 2.1, 3.2, BG_CARD2)
        add_text_box(slide, x, 3.9, 2.1, 0.6, icon,
                     font_size=36, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, 4.6, 2.1, 0.4, title,
                     font_size=14, bold=True, color=GREEN_SOFT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 5.1, 1.9, 1.5, desc,
                     font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_footer(slide)

    # --- スライド3: AI病害虫診断 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '🔬 AI病害虫診断',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    # 病害
    add_shape_bg(slide, 0.5, 1.2, 4.3, 3.0, BG_CARD)
    add_text_box(slide, 0.7, 1.3, 4, 0.4, '🍂 対応病害（5種）',
                 font_size=16, bold=True, color=ACCENT_AMBER)
    add_bullet_list(slide, 0.7, 1.8, 4, 2.3, [
        '・ いもち病（葉いもち・穂いもち・首いもち）',
        '・ 紋枯病',
        '・ 白葉枯病',
        '・ 縞葉枯病',
        '・ 細菌性穀枯病',
    ], font_size=13, color=TEXT_WHITE)

    # 害虫
    add_shape_bg(slide, 5.2, 1.2, 4.3, 3.0, BG_CARD)
    add_text_box(slide, 5.4, 1.3, 4, 0.4, '🦗 対応害虫（6種）',
                 font_size=16, bold=True, color=ACCENT_AMBER)
    add_bullet_list(slide, 5.4, 1.8, 4, 2.3, [
        '・ トビイロウンカ',
        '・ セジロウンカ',
        '・ イネカメムシ',
        '・ クモヘリカメムシ',
        '・ ニカメイチュウ',
        '・ イネドロオイムシ',
    ], font_size=13, color=TEXT_WHITE)

    # 出力項目
    add_shape_bg(slide, 0.5, 4.5, 9, 2.5, BG_CARD2)
    add_text_box(slide, 0.7, 4.6, 8.5, 0.4, '📊 解析出力',
                 font_size=16, bold=True, color=GREEN_SOFT)
    outputs = [
        '✅ 病害虫名（学名付き）      ✅ 確信度（%）         ✅ 被害面積率',
        '✅ 重度判定（軽度〜重度）    ✅ 推奨農薬             ✅ 緊急度（即時/1週間/経過観察）',
        '✅ 視覚的根拠の説明          ✅ 鑑別診断             ✅ 具体的な対応策の提示',
    ]
    add_bullet_list(slide, 0.7, 5.1, 8.5, 2, outputs, font_size=12, color=TEXT_WHITE)
    add_footer(slide)

    # --- スライド4: 画面紹介（ダッシュボード＆アップロード）---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '📱 画面紹介',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_text_box(slide, 0.5, 1.1, 4.5, 0.3, 'ダッシュボード',
                 font_size=14, bold=True, color=GREEN_SOFT)
    add_image_safe(slide, 'screenshot_dashboard.png', 0.5, 1.5, width=4.4)

    add_text_box(slide, 5.2, 1.1, 4.5, 0.3, '画像アップロード',
                 font_size=14, bold=True, color=GREEN_SOFT)
    add_image_safe(slide, 'screenshot_upload.png', 5.2, 1.5, width=4.4)
    add_footer(slide)

    # --- スライド5: 画面紹介（レポート＆圃場管理）---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '📱 画面紹介',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_text_box(slide, 0.5, 1.1, 4.5, 0.3, '解析レポート（5タブ）',
                 font_size=14, bold=True, color=GREEN_SOFT)
    add_image_safe(slide, 'screenshot_report.png', 0.5, 1.5, width=4.4)

    add_text_box(slide, 5.2, 1.1, 4.5, 0.3, '圃場管理（衛星写真マップ）',
                 font_size=14, bold=True, color=GREEN_SOFT)
    add_image_safe(slide, 'screenshot_fields.png', 5.2, 1.5, width=4.4)
    add_footer(slide)

    # --- スライド6: 圃場ポリゴン描画 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 8, 0.5, '📐 ポリゴン描画 × 面積自動算出',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_polygon.png', 0.5, 1.2, width=5.5)

    add_shape_bg(slide, 6.3, 1.2, 3.3, 5.5, BG_CARD)
    add_text_box(slide, 6.5, 1.3, 3, 0.4, '操作手順',
                 font_size=16, bold=True, color=GREEN_SOFT)
    steps = [
        '❶ 衛星写真マップを表示',
        '',
        '❷ 圃場の角を順番にクリック',
        '   （3点以上で多角形を描画）',
        '',
        '❸ 面積を自動算出',
        '   ・ ha（ヘクタール）',
        '   ・ m²（平方メートル）',
        '   ・ 反（たん）',
        '',
        '❹ 中心座標を自動取得',
        '',
        '❺ 「圃場を登録」で保存',
    ]
    add_bullet_list(slide, 6.5, 1.8, 3, 4.5, steps, font_size=12, color=TEXT_WHITE)
    add_footer(slide)

    # --- スライド7: 技術構成 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '⚙️ 技術構成',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    layers = [
        ('フロントエンド', 'HTML / CSS / JavaScript\nLeaflet.js（衛星写真マップ）', GREEN_SOFT),
        ('バックエンド', 'Python / Flask\nREST API', ACCENT_CYAN),
        ('AI解析エンジン', 'Claude Vision API（Anthropic）\n5モジュール解析プロンプト', ACCENT_AMBER),
        ('データベース', 'SQLite\n圃場情報・解析結果の永続化', TEXT_MUTED),
        ('地図タイル', 'Esri World Imagery（無料）\n衛星写真表示', GREEN_LIGHT),
    ]
    for i, (title, desc, color) in enumerate(layers):
        y = 1.2 + i * 1.15
        add_shape_bg(slide, 0.5, y, 9, 1.0, BG_CARD)
        add_text_box(slide, 0.7, y + 0.05, 2.5, 0.9, title,
                     font_size=16, bold=True, color=color)
        add_text_box(slide, 3.5, y + 0.1, 5.8, 0.8, desc,
                     font_size=13, color=TEXT_WHITE)
    add_footer(slide)

    # --- スライド8: コスト ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '💰 コスト目安',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    costs = [
        ('ソフトウェア', '無料', '（オープンソース）'),
        ('Claude API', '約 ¥5〜15/回', '（画像1枚あたり）'),
        ('地図タイル', '無料', '（Esri World Imagery）'),
        ('サーバー費用', '無料', '（ローカル動作）'),
    ]
    for i, (item, cost, note) in enumerate(costs):
        y = 1.3 + i * 0.9
        add_shape_bg(slide, 0.5, y, 9, 0.75, BG_CARD)
        add_text_box(slide, 0.7, y + 0.1, 3, 0.5, item,
                     font_size=15, bold=True, color=TEXT_WHITE)
        add_text_box(slide, 4.0, y + 0.1, 2.5, 0.5, cost,
                     font_size=18, bold=True, color=GREEN_LIGHT)
        add_text_box(slide, 6.5, y + 0.15, 3, 0.4, note,
                     font_size=12, color=TEXT_MUTED)

    add_shape_bg(slide, 2, 5.2, 6, 1.3, RGBColor(0x1B, 0x3A, 0x26))
    add_text_box(slide, 2.2, 5.3, 5.5, 0.4, '月間100回解析の場合',
                 font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2.2, 5.7, 5.5, 0.6, '約 ¥500〜1,500 / 月',
                 font_size=32, bold=True, color=GREEN_LIGHT, alignment=PP_ALIGN.CENTER)
    add_footer(slide)

    # --- スライド9: 今後の拡張 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '🚀 今後の拡張予定',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    roadmap = [
        ('📡', 'RTK基準局データ連携', '±2cm精度の圃場境界特定'),
        ('📈', '時系列変化トラッキング', '過去画像との比較で推移を把握'),
        ('📱', 'LINEアラート通知', '緊急検出時の自動通知'),
        ('👥', 'マルチユーザー対応', '複数農家の圃場を一元管理'),
        ('🌤️', '気象データ連携', '気象庁APIとの統合で予測精度向上'),
        ('🤖', 'ドローン散布計画の自動生成', 'AIが最適散布ルートを提案'),
    ]
    for i, (icon, title, desc) in enumerate(roadmap):
        y = 1.2 + i * 0.95
        add_shape_bg(slide, 0.5, y, 9, 0.8, BG_CARD)
        add_text_box(slide, 0.7, y + 0.1, 0.5, 0.5, icon, font_size=20)
        add_text_box(slide, 1.3, y + 0.1, 3.5, 0.5, title,
                     font_size=15, bold=True, color=TEXT_WHITE)
        add_text_box(slide, 5.0, y + 0.15, 4.5, 0.4, desc,
                     font_size=13, color=TEXT_MUTED)
    add_footer(slide)

    out_path = os.path.join(DOCS_DIR, 'AgriScan_Pro_概要資料.pptx')
    prs.save(out_path)
    print(f'✅ 概要資料を保存: {out_path}')
    return out_path


# =============================================
# 操作マニュアル
# =============================================
def create_manual():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- 表紙 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(3.2), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_text_box(slide, 1, 1.5, 8, 1, '📖', font_size=60, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.2, 8, 0.8, 'AgriScan Pro',
                 font_size=40, bold=True, color=GREEN_LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.5, 8, 0.6, '操作マニュアル',
                 font_size=28, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 8, 0.4, 'セットアップ ・ 各画面の使い方 ・ トラブルシューティング',
                 font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6.0, 8, 0.4, 'v1.0  |  2026年3月',
                 font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # --- 目次 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '📋 目次',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    toc = [
        ('01', '起動方法', 'セットアップ・APIキー設定'),
        ('02', 'ダッシュボード', '全体の健康状態を一覧'),
        ('03', '画像アップロード', 'ドローン画像のAI解析'),
        ('04', '解析レポートの見方', '5タブの詳細結果'),
        ('05', '圃場管理', '衛星マップ・圃場一覧'),
        ('06', '圃場の追加', 'ポリゴン描画・面積自動算出'),
        ('07', 'トラブルシューティング', 'よくある問題と対処法'),
    ]
    for i, (num, title, desc) in enumerate(toc):
        y = 1.3 + i * 0.75
        add_shape_bg(slide, 0.5, y, 9, 0.6, BG_CARD)
        add_text_box(slide, 0.7, y + 0.05, 0.6, 0.4, num,
                     font_size=18, bold=True, color=GREEN_PRIMARY)
        add_text_box(slide, 1.5, y + 0.05, 3.5, 0.4, title,
                     font_size=16, bold=True, color=TEXT_WHITE)
        add_text_box(slide, 5.5, y + 0.1, 4, 0.3, desc,
                     font_size=12, color=TEXT_MUTED)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- 起動方法 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '01  起動方法',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_shape_bg(slide, 0.5, 1.2, 4.4, 5.5, BG_CARD)
    add_text_box(slide, 0.7, 1.3, 4, 0.4, '🔧 初回セットアップ',
                 font_size=16, bold=True, color=ACCENT_AMBER)
    add_bullet_list(slide, 0.7, 1.8, 4, 3, [
        '❶ プロジェクトフォルダに移動',
        '   cd ドローン農業関連/agriscan-pro',
        '',
        '❷ Python仮想環境を作成',
        '   python3 -m venv backend/venv',
        '',
        '❸ 仮想環境を有効化',
        '   source backend/venv/bin/activate',
        '',
        '❹ パッケージをインストール',
        '   pip install -r backend/requirements.txt',
    ], font_size=11, color=TEXT_WHITE)

    add_shape_bg(slide, 5.2, 1.2, 4.4, 3.0, BG_CARD)
    add_text_box(slide, 5.4, 1.3, 4, 0.4, '🚀 毎回の起動',
                 font_size=16, bold=True, color=GREEN_SOFT)
    add_bullet_list(slide, 5.4, 1.8, 4, 2, [
        '❶ source backend/venv/bin/activate',
        '',
        '❷ python backend/app.py',
        '',
        '❸ ブラウザで開く:',
        '   http://localhost:5001/',
    ], font_size=12, color=TEXT_WHITE)

    add_shape_bg(slide, 5.2, 4.5, 4.4, 2.2, BG_CARD2)
    add_text_box(slide, 5.4, 4.6, 4, 0.4, '🔑 AIモードを有効にする',
                 font_size=14, bold=True, color=ACCENT_CYAN)
    add_bullet_list(slide, 5.4, 5.1, 4, 1.2, [
        'export ANTHROPIC_API_KEY=sk-ant-...',
        'python backend/app.py',
        '',
        '※ キー未設定 → デモモードで動作',
    ], font_size=11, color=TEXT_WHITE)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- ダッシュボード ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '02  ダッシュボード',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_dashboard.png', 0.5, 1.2, width=5.5)

    add_shape_bg(slide, 6.3, 1.2, 3.3, 5.5, BG_CARD)
    add_text_box(slide, 6.5, 1.3, 3, 0.4, '表示項目',
                 font_size=16, bold=True, color=GREEN_SOFT)
    items = [
        '📊 管理圃場数',
        '📐 総管理面積（ha）',
        '⚠️ 直近14日のアラート',
        '💯 前回解析スコア',
        '',
        '🎯 総合健康スコアリング',
        '   （円形ゲージ表示）',
        '',
        '🚨 緊急アラート一覧',
        '   （即時対応が必要な項目）',
        '',
        '📈 モジュール別スコア',
        '   病害虫 / 雑草 / インフラ / 散布',
    ]
    add_bullet_list(slide, 6.5, 1.8, 3, 4.5, items, font_size=11, color=TEXT_WHITE)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- アップロード ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 8, 0.5, '03  画像アップロード・AI解析',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_upload.png', 0.5, 1.2, width=5.0)

    add_shape_bg(slide, 5.8, 1.2, 3.8, 3.5, BG_CARD)
    add_text_box(slide, 6.0, 1.3, 3.5, 0.4, '手順',
                 font_size=16, bold=True, color=GREEN_SOFT)
    add_bullet_list(slide, 6.0, 1.8, 3.5, 3, [
        '❶ 画像を選択 or ドラッグ&ドロップ',
        '   対応: JPG, PNG, TIFF, WebP',
        '   最大: 50MB',
        '',
        '❷ 撮影情報を入力',
        '   圃場・品種・ステージ・高度',
        '',
        '❸「AI解析を開始」をクリック',
        '',
        '❹ 約10-30秒で結果表示',
    ], font_size=11, color=TEXT_WHITE)

    add_shape_bg(slide, 5.8, 5.0, 3.8, 1.8, BG_CARD2)
    add_text_box(slide, 6.0, 5.1, 3.5, 0.3, '📸 推奨撮影条件',
                 font_size=14, bold=True, color=ACCENT_CYAN)
    add_bullet_list(slide, 6.0, 5.5, 3.5, 1, [
        '高度: 15〜30m',
        '時間帯: 午前9時〜午後3時',
        '天候: 曇り〜晴れ（雨天は避ける）',
    ], font_size=11, color=TEXT_WHITE)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- レポート ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 8, 0.5, '04  解析レポートの見方',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_report.png', 0.5, 1.2, width=5.0)

    tabs = [
        ('🔬', '病害虫', '検出結果・確信度・推奨農薬', GREEN_LIGHT),
        ('🌿', '雑草密度', 'ヒートマップ・被覆率', GREEN_SOFT),
        ('🏗️', 'インフラ', '畦畔崩れ・水路詰まり', ACCENT_CYAN),
        ('✈️', '散布記録', '散布履歴タイムライン', ACCENT_AMBER),
        ('📊', '総合スコア', '0-100評価・次アクション', WHITE),
    ]
    for i, (icon, name, desc, color) in enumerate(tabs):
        y = 1.2 + i * 1.1
        add_shape_bg(slide, 5.8, y, 3.8, 0.95, BG_CARD)
        add_text_box(slide, 6.0, y + 0.05, 0.4, 0.4, icon, font_size=18)
        add_text_box(slide, 6.4, y + 0.05, 1.5, 0.4, name,
                     font_size=15, bold=True, color=color)
        add_text_box(slide, 6.0, y + 0.5, 3.5, 0.3, desc,
                     font_size=11, color=TEXT_MUTED)

    # スコア判定
    add_shape_bg(slide, 5.8, 6.8, 3.8, 0.5, BG_CARD2)
    add_text_box(slide, 6.0, 6.85, 3.5, 0.3,
                 '🟢80-100 良好  🟡60-79 注意  🟠40-59 要対処  🔴0-39 緊急',
                 font_size=9, color=TEXT_MUTED)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- 圃場管理 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.5, '05  圃場管理',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_fields.png', 0.5, 1.2, width=5.5)

    add_shape_bg(slide, 6.3, 1.2, 3.3, 5.5, BG_CARD)
    add_text_box(slide, 6.5, 1.3, 3, 0.4, '画面構成',
                 font_size=16, bold=True, color=GREEN_SOFT)
    add_bullet_list(slide, 6.5, 1.8, 3, 4.5, [
        '🛰️ 衛星写真マップ',
        '   Esri World Imagery（無料）',
        '',
        '📍 圃場マーカー',
        '   🟢 良好  🟡 注意  🔴 要対処',
        '   クリックで詳細ポップアップ',
        '',
        '📊 統計カード',
        '   圃場数/面積/ステータス',
        '',
        '🃏 圃場カード一覧',
        '   名前・農家・面積・品種・',
        '   健康スコア・ステータス',
    ], font_size=11, color=TEXT_WHITE)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- 圃場の追加 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 8, 0.5, '06  圃場の追加（ポリゴン描画）',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    add_image_safe(slide, 'screenshot_add_field.png', 0.3, 1.2, width=4.5)
    add_image_safe(slide, 'screenshot_polygon.png', 5.0, 1.2, width=4.7)

    add_shape_bg(slide, 0.5, 5.5, 9, 1.8, BG_CARD)
    add_text_box(slide, 0.7, 5.6, 8.5, 0.3, '操作手順',
                 font_size=14, bold=True, color=GREEN_SOFT)
    add_text_box(slide, 0.7, 6.0, 8.5, 1.0,
                 '❶「➕ 圃場を追加」をクリック  →  '
                 '❷ 圃場名・農家名・品種を入力  →  '
                 '❸ 衛星写真上で田んぼの角を順番にクリック（3点以上）\n'
                 '❹ 面積が自動算出（ha/m²/反）  →  '
                 '❺「↩ 1点戻す」で修正可能  →  '
                 '❻「🌾 圃場を登録」で保存完了',
                 font_size=13, color=TEXT_WHITE)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    # --- トラブルシューティング ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.6, 0.3, 8, 0.5, '07  トラブルシューティング',
                 font_size=28, bold=True, color=GREEN_LIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.9), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_PRIMARY
    line.line.fill.background()

    troubles = [
        ('サーバーが起動しない', 'ポート5001が使用中', 'lsof -i :5001 で確認し kill -9 で停止'),
        ('Claude API 無効と表示', 'APIキー未設定', 'export ANTHROPIC_API_KEY=... を実行後に再起動'),
        ('画像アップロードエラー', '非対応形式 or 50MB超過', 'JPG/PNG/TIFF形式で50MB以下にする'),
        ('マップが表示されない', 'インターネット未接続', 'WiFi/LAN接続を確認'),
        ('デモモードから抜けない', 'サーバー再起動が必要', 'Ctrl+C → APIキー設定 → 再起動'),
    ]
    for i, (problem, cause, solution) in enumerate(troubles):
        y = 1.2 + i * 1.2
        add_shape_bg(slide, 0.5, y, 9, 1.05, BG_CARD)
        add_text_box(slide, 0.7, y + 0.05, 2.5, 0.4, f'⚠️ {problem}',
                     font_size=13, bold=True, color=ACCENT_AMBER)
        add_text_box(slide, 0.7, y + 0.5, 2.5, 0.4, f'原因: {cause}',
                     font_size=11, color=TEXT_MUTED)
        add_text_box(slide, 4.0, y + 0.15, 5.3, 0.7, f'💡 {solution}',
                     font_size=13, color=GREEN_SOFT)
    add_footer(slide, "AgriScan Pro 操作マニュアル v1.0")

    out_path = os.path.join(DOCS_DIR, 'AgriScan_Pro_操作マニュアル.pptx')
    prs.save(out_path)
    print(f'✅ 操作マニュアルを保存: {out_path}')
    return out_path


if __name__ == '__main__':
    create_overview()
    create_manual()
    print('\n🎉 2つのパワーポイント資料を生成しました！')
